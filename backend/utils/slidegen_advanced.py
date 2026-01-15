"""
Advanced Slide Generation Utility
Handles intelligent content analysis, adaptive chunking, template selection, and PPTX/PDF generation
"""

import re
import os
import subprocess
from typing import List, Dict, Tuple, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from io import BytesIO
import tempfile

# Template definitions
TEMPLATES = {
    "academic": {
        "primary_color": RGBColor(25, 55, 109),      # Navy blue
        "accent_color": RGBColor(70, 130, 180),      # Steel blue
        "background_color": RGBColor(245, 248, 250), # Light blue-gray
        "font_family": "Calibri",
        "name": "Academic"
    },
    "business": {
        "primary_color": RGBColor(0, 51, 102),       # Dark blue
        "accent_color": RGBColor(204, 51, 0),        # Dark orange
        "background_color": RGBColor(255, 255, 255), # White
        "font_family": "Arial",
        "name": "Business"
    },
    "creative": {
        "primary_color": RGBColor(102, 51, 153),     # Purple
        "accent_color": RGBColor(255, 102, 178),     # Pink
        "background_color": RGBColor(245, 240, 245), # Light purple
        "font_family": "Calibri",
        "name": "Creative"
    },
    "research": {
        "primary_color": RGBColor(34, 102, 68),      # Forest green
        "accent_color": RGBColor(100, 150, 100),     # Light green
        "background_color": RGBColor(240, 248, 245), # Mint
        "font_family": "Calibri",
        "name": "Research"
    },
    "default": {
        "primary_color": RGBColor(31, 78, 121),      # Navy
        "accent_color": RGBColor(79, 129, 189),      # Light blue
        "background_color": RGBColor(255, 255, 255), # White
        "font_family": "Calibri",
        "name": "Default"
    }
}


class ContentAnalyzer:
    """Analyzes content to determine structure and properties"""
    
    @staticmethod
    def analyze(text: str) -> Dict:
        """
        Analyze content and return metadata
        Returns: {
            'word_count': int,
            'section_count': int,
            'content_type': str (academic|business|creative|research),
            'density': str (low|medium|high),
            'estimated_slides': int,
            'sections': List[Dict]
        }
        """
        word_count = len(text.split())
        
        # Detect sections
        sections = ContentAnalyzer._extract_sections(text)
        section_count = len(sections)
        
        # Detect content type by keywords
        content_type = ContentAnalyzer._detect_content_type(text)
        
        # Determine content density
        density = "low" if word_count < 500 else "medium" if word_count < 2000 else "high"
        
        # Estimate slide count
        estimated_slides = ContentAnalyzer._estimate_slide_count(word_count, section_count, density)
        
        return {
            'word_count': word_count,
            'section_count': section_count,
            'content_type': content_type,
            'density': density,
            'estimated_slides': estimated_slides,
            'sections': sections
        }
    
    @staticmethod
    def _extract_sections(text: str) -> List[Dict]:
        """Extract natural sections from text"""
        sections = []
        
        # Split by common headers (##, **, etc.)
        # Also split by double newlines (paragraph breaks)
        parts = re.split(r'\n\n+|^#+\s+', text, flags=re.MULTILINE)
        
        for part in parts:
            part = part.strip()
            if len(part) > 20:  # Only sections with meaningful content
                # First line as title, rest as content
                lines = part.split('\n')
                title = lines[0][:100] if lines else "Section"
                content = '\n'.join(lines[1:]) if len(lines) > 1 else part
                
                sections.append({
                    'title': title,
                    'content': content,
                    'word_count': len(content.split())
                })
        
        return sections[:50]  # Cap at 50 sections for performance
    
    @staticmethod
    def _detect_content_type(text: str) -> str:
        """Detect content type by analyzing keywords"""
        text_lower = text.lower()
        
        # Academic indicators
        academic_keywords = ['research', 'study', 'hypothesis', 'methodology', 'conclusion', 
                           'abstract', 'literature', 'analysis', 'experiment', 'theory']
        
        # Business indicators
        business_keywords = ['revenue', 'profit', 'market', 'strategy', 'quarterly', 
                           'performance', 'growth', 'stakeholder', 'client', 'investment']
        
        # Research indicators
        research_keywords = ['data', 'analysis', 'findings', 'methodology', 'quantitative',
                           'qualitative', 'variables', 'statistical', 'correlation']
        
        academic_score = sum(1 for kw in academic_keywords if kw in text_lower)
        business_score = sum(1 for kw in business_keywords if kw in text_lower)
        research_score = sum(1 for kw in research_keywords if kw in text_lower)
        
        scores = {
            'academic': academic_score,
            'business': business_score,
            'research': research_score
        }
        
        detected_type = max(scores, key=scores.get) if max(scores.values()) > 0 else 'default'
        return detected_type
    
    @staticmethod
    def _estimate_slide_count(word_count: int, section_count: int, density: str) -> int:
        """Intelligently estimate slide count"""
        if word_count < 500:
            return 3 + section_count
        elif word_count < 2000:
            return 5 + min(section_count, 10)
        else:
            # For large content, cap at 15 slides (rest will be summarized)
            return min(10 + section_count // 2, 15)


class SlideChunker:
    """Intelligently chunks content for slides"""
    
    @staticmethod
    def chunk_content(sections: List[Dict], estimated_slides: int, max_slides: int = 50) -> List[Dict]:
        """
        Create logical slide chunks from sections
        Respects max_slides limit and creates balanced chunks
        """
        if estimated_slides > max_slides:
            # Summarize or compress content
            sections = SlideChunker._compress_sections(sections, max_slides)
        
        chunks = []
        
        for section in sections:
            title = section['title']
            content = section['content']
            word_count = section['word_count']
            
            # Determine how many slides this section needs
            if word_count < 100:
                slides_needed = 1
            elif word_count < 300:
                slides_needed = 1
            else:
                slides_needed = min((word_count // 200) + 1, 3)
            
            # Split content into bullet points
            bullet_points = SlideChunker._extract_bullet_points(content, slides_needed)
            
            # Create slide chunks
            for i, points in enumerate(bullet_points):
                chunks.append({
                    'title': title if i == 0 else f"{title} (cont.)",
                    'bullet_points': points,
                    'type': 'content'
                })
        
        return chunks[:max_slides]
    
    @staticmethod
    def _compress_sections(sections: List[Dict], target_slides: int) -> List[Dict]:
        """Compress sections to fit target slide count"""
        if not sections:
            return sections
        
        # Merge smaller sections
        compressed = []
        current_merged = None
        
        for section in sections:
            if current_merged is None:
                current_merged = section
            elif section['word_count'] < 100 and current_merged['word_count'] < 200:
                # Merge small sections
                current_merged['content'] += '\n' + section['content']
                current_merged['word_count'] += section['word_count']
            else:
                compressed.append(current_merged)
                current_merged = section
        
        if current_merged:
            compressed.append(current_merged)
        
        return compressed
    
    @staticmethod
    def _extract_bullet_points(content: str, num_slides: int) -> List[List[str]]:
        """Extract bullet points from content"""
        # Split by sentences or lines
        sentences = re.split(r'(?<=[.!?])\s+|\n', content.strip())
        sentences = [s.strip() for s in sentences if len(s.strip()) > 10]
        
        if not sentences:
            return [['No content available']]
        
        # Distribute sentences across slides
        points_per_slide = max(2, len(sentences) // num_slides)
        slides = []
        
        for i in range(num_slides):
            start = i * points_per_slide
            end = start + points_per_slide if i < num_slides - 1 else len(sentences)
            slide_points = sentences[start:end]
            
            # Limit to 5 bullet points per slide
            slides.append(slide_points[:5])
        
        return [s for s in slides if s]  # Remove empty slides


class SlideBuilder:
    """Builds PPTX presentation"""
    
    @staticmethod
    def build_presentation(
        title: str,
        chunks: List[Dict],
        template: str = "default",
        subtitle: Optional[str] = None
    ) -> Presentation:
        """Build PPTX presentation from chunks"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        template_config = TEMPLATES.get(template, TEMPLATES["default"])
        
        # Title slide
        SlideBuilder._add_title_slide(prs, title, subtitle or "Generated Presentation", template_config)
        
        # Content slides
        for chunk in chunks:
            if chunk['type'] == 'content':
                SlideBuilder._add_content_slide(
                    prs,
                    chunk['title'],
                    chunk['bullet_points'],
                    template_config
                )
        
        # Summary slide
        SlideBuilder._add_summary_slide(prs, title, len(chunks), template_config)
        
        return prs
    
    @staticmethod
    def _add_title_slide(prs: Presentation, title: str, subtitle: str, template: Dict):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = template['background_color']
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        title_p.font.color.rgb = template['primary_color']
        title_p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = subtitle
        subtitle_p.font.size = Pt(24)
        subtitle_p.font.color.rgb = template['accent_color']
        subtitle_p.alignment = PP_ALIGN.CENTER
    
    @staticmethod
    def _add_content_slide(prs: Presentation, title: str, bullet_points: List[str], template: Dict):
        """Add content slide with bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = template['background_color']
        
        # Title bar
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = template['primary_color']
        title_shape.line.color.rgb = template['primary_color']
        
        # Title text
        title_frame = title_shape.text_frame
        title_frame.margin_left = Inches(0.3)
        title_frame.margin_top = Inches(0.1)
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Content area
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_before = Pt(12)
            p.space_after = Pt(12)
    
    @staticmethod
    def _add_summary_slide(prs: Presentation, title: str, slide_count: int, template: Dict):
        """Add summary slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = template['background_color']
        
        # Title bar
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = template['primary_color']
        title_shape.line.color.rgb = template['primary_color']
        
        title_frame = title_shape.text_frame
        title_frame.margin_left = Inches(0.3)
        title_frame.margin_top = Inches(0.1)
        title_p = title_frame.paragraphs[0]
        title_p.text = "Summary"
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(3))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        summary_text = f"""
Thank you for reviewing this presentation!

Presentation: {title}
Total Slides: {slide_count + 2}
Generated with: Martian AI
        """.strip()
        
        p = text_frame.paragraphs[0]
        p.text = summary_text
        p.font.size = Pt(22)
        p.font.color.rgb = template['primary_color']
        p.alignment = PP_ALIGN.CENTER


class PDFConverter:
    """Converts PPTX to PDF"""
    
    @staticmethod
    def convert_pptx_to_pdf(pptx_bytes: BytesIO, output_path: str) -> bool:
        """
        Convert PPTX to PDF
        Tries LibreOffice first, falls back to reporting if available
        Returns: True if successful, False otherwise
        """
        # Method 1: Try LibreOffice (best quality)
        try:
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                tmp.write(pptx_bytes.getvalue())
                tmp_path = tmp.name
            
            try:
                subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', 
                     os.path.dirname(output_path), tmp_path],
                    check=True,
                    capture_output=True,
                    timeout=30
                )
                # Move converted file to target location
                converted_path = tmp_path.replace('.pptx', '.pdf')
                if os.path.exists(converted_path):
                    os.rename(converted_path, output_path)
                    return True
            finally:
                if os.path.exists(tmp_path):
                    os.remove(tmp_path)
        except (FileNotFoundError, subprocess.TimeoutExpired, subprocess.CalledProcessError):
            pass
        
        # Method 2: Fallback - return False, let client handle it
        return False


class SlideGenerator:
    """Main slide generation orchestrator"""
    
    @staticmethod
    def generate_from_text(
        text: str,
        title: str,
        prompt: Optional[str] = None,
        template: Optional[str] = None,
        include_pdf: bool = False
    ) -> Dict:
        """
        Generate slides from text
        Returns: {'pptx': BytesIO, 'pdf': BytesIO|None, 'metadata': {...}}
        """
        # Validate input
        if len(text) == 0:
            raise ValueError("Text content is required")
        
        # Analyze content
        analysis = ContentAnalyzer.analyze(text)
        
        # Auto-detect template if not provided
        detected_template = template or analysis['content_type']
        if detected_template not in TEMPLATES:
            detected_template = "default"
        
        # Intelligent template override from prompt
        if prompt:
            detected_template = SlideGenerator._detect_template_from_prompt(prompt, detected_template)
        
        # Chunk content
        chunks = SlideChunker.chunk_content(
            analysis['sections'],
            analysis['estimated_slides']
        )
        
        if not chunks:
            raise ValueError("No valid content to generate slides from")
        
        # Build PPTX
        prs = SlideBuilder.build_presentation(
            title=title,
            chunks=chunks,
            template=detected_template,
            subtitle=prompt or "Generated Presentation"
        )
        
        # Save PPTX to bytes
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        
        # Attempt PDF conversion if requested
        pdf_io = None
        pdf_success = False
        if include_pdf:
            try:
                with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                    pdf_path = tmp.name
                
                pdf_success = PDFConverter.convert_pptx_to_pdf(pptx_io, pdf_path)
                
                if pdf_success and os.path.exists(pdf_path):
                    with open(pdf_path, 'rb') as f:
                        pdf_io = BytesIO(f.read())
                    os.remove(pdf_path)
            except Exception as e:
                print(f"PDF conversion failed: {e}")
        
        return {
            'pptx': pptx_io,
            'pdf': pdf_io,
            'metadata': {
                'title': title,
                'word_count': analysis['word_count'],
                'slide_count': len(chunks) + 2,  # +2 for title and summary
                'template': detected_template,
                'content_type': analysis['content_type'],
                'pdf_available': pdf_success
            }
        }
    
    @staticmethod
    def _detect_template_from_prompt(prompt: str, default_template: str) -> str:
        """Detect template preference from user prompt"""
        prompt_lower = prompt.lower()
        
        if any(word in prompt_lower for word in ['academic', 'research', 'formal']):
            return 'academic'
        elif any(word in prompt_lower for word in ['business', 'corporate', 'professional']):
            return 'business'
        elif any(word in prompt_lower for word in ['creative', 'colorful', 'modern']):
            return 'creative'
        elif any(word in prompt_lower for word in ['data', 'analysis', 'science']):
            return 'research'
        
        return default_template


# Main entry point for external use
def generate_slides(
    content: str,
    title: str,
    prompt: Optional[str] = None,
    template: Optional[str] = None,
    export_format: str = "pptx"  # "pptx", "pdf", or "both"
) -> Dict:
    """
    Public API for slide generation
    
    Args:
        content: Text content for slides
        title: Presentation title
        prompt: User's custom instructions/template preference
        template: Force specific template (academic, business, creative, research, default)
        export_format: "pptx", "pdf", or "both"
    
    Returns:
        {'pptx': BytesIO, 'pdf': BytesIO|None, 'metadata': {...}}
    """
    include_pdf = export_format in ['pdf', 'both']
    
    return SlideGenerator.generate_from_text(
        text=content,
        title=title,
        prompt=prompt,
        template=template,
        include_pdf=include_pdf
    )
